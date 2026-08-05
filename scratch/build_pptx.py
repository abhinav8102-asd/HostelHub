import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_slide_layout = prs.slide_layouts[6]

    # Color Palette
    DEEP_MAROON = RGBColor(138, 13, 36)    # #8A0D24
    CRIMSON = RGBColor(179, 16, 49)       # #B31031
    DARK_PLUM = RGBColor(76, 6, 21)        # #4C0615
    TEXT_DARK = RGBColor(15, 23, 42)       # #0F172A
    TEXT_MUTED = RGBColor(100, 116, 139)   # #64748B
    CARD_BG = RGBColor(250, 248, 248)      # #FAF8F8
    BORDER_GRAY = RGBColor(226, 232, 240)  # #E2E8F0
    WHITE = RGBColor(255, 255, 255)
    GOLD = RGBColor(217, 119, 6)
    GREEN = RGBColor(22, 101, 52)

    img_dir = r"C:\Users\abhin\.gemini\antigravity-ide\brain\a135e2ae-3ec5-4446-a8af-46a5537b2c14"
    img_attendance = os.path.join(img_dir, "media__1785667420651.png")
    img_complaint = os.path.join(img_dir, "media__1785587403608.png")
    img_profile = os.path.join(img_dir, "media__1785587242838.png")
    img_devs = os.path.join(img_dir, "media__1785586595628.png")

    def add_header(slide, title_text, category_text="HOSTELHUB PLATFORM"):
        # Top banner background
        top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.1))
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = DEEP_MAROON
        top_bar.line.fill.background()

        # Gold accent stripe
        stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.1), Inches(13.333), Inches(0.06))
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = GOLD
        stripe.line.fill.background()

        # Logo Box
        logo_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(0.2), Inches(0.7), Inches(0.7))
        logo_box.fill.solid()
        logo_box.fill.fore_color.rgb = WHITE
        logo_box.line.fill.background()
        tf = logo_box.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = "H"
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = DEEP_MAROON
        p.alignment = PP_ALIGN.CENTER

        # Category & Title Text
        tx_box = slide.shapes.add_textbox(Inches(1.4), Inches(0.15), Inches(11.0), Inches(0.8))
        tf = tx_box.text_frame
        p0 = tf.paragraphs[0]
        p0.text = category_text.upper()
        p0.font.size = Pt(9)
        p0.font.bold = True
        p0.font.color.rgb = GOLD

        p1 = tf.add_paragraph()
        p1.text = title_text
        p1.font.size = Pt(22)
        p1.font.bold = True
        p1.font.color.rgb = WHITE

    # ==========================================
    # SLIDE 1: COVER SLIDE
    # ==========================================
    slide1 = prs.slides.add_slide(blank_slide_layout)
    bg1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = DEEP_MAROON
    bg1.line.fill.background()

    # Decorative shapes
    shape1 = slide1.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, Inches(7.0), Inches(0), Inches(6.333), Inches(7.5))
    shape1.fill.solid()
    shape1.fill.fore_color.rgb = DARK_PLUM
    shape1.line.fill.background()
    shape1.rotation = 180

    # Logo
    logo1 = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(1.5), Inches(1.2), Inches(1.2))
    logo1.fill.solid()
    logo1.fill.fore_color.rgb = WHITE
    logo1.line.fill.background()
    tf = logo1.text_frame
    p = tf.paragraphs[0]
    p.text = "H"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = DEEP_MAROON
    p.alignment = PP_ALIGN.CENTER

    tx = slide1.shapes.add_textbox(Inches(1.0), Inches(3.0), Inches(11.0), Inches(3.5))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.text = "HostelHub"
    p.font.size = Pt(46)
    p.font.bold = True
    p.font.color.rgb = WHITE

    p2 = tf.add_paragraph()
    p2.text = "All-in-One Digital Hostel Management Platform"
    p2.font.size = Pt(24)
    p2.font.bold = True
    p2.font.color.rgb = GOLD

    p3 = tf.add_paragraph()
    p3.text = "Streamlining Complaints, Daily Attendance, Mess Regulations & Real-Time Communication"
    p3.font.size = Pt(14)
    p3.font.color.rgb = WHITE

    p4 = tf.add_paragraph()
    p4.text = "\nDeveloped by Abhinav Kumar (Lead Full-Stack Developer) | August 2026"
    p4.font.size = Pt(12)
    p4.font.color.rgb = RGBColor(226, 232, 240)

    # ==========================================
    # SLIDE 2: PROBLEM STATEMENT & VISION
    # ==========================================
    slide2 = prs.slides.add_slide(blank_slide_layout)
    add_header(slide2, "Problem Statement & Digital Vision")

    # Card 1: Traditional Challenges
    card1 = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.3))
    card1.fill.solid()
    card1.fill.fore_color.rgb = CARD_BG
    card1.line.color.rgb = BORDER_GRAY
    tf1 = card1.text_frame
    tf1.word_wrap = True
    p = tf1.paragraphs[0]
    p.text = "⚠️ Traditional Hostel Challenges"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = DEEP_MAROON

    bullets1 = [
        "Manual Paper Register Attendance leads to proxy errors and lost records.",
        "Unorganized Maintenance Complaints cause long delays in repair dispatch.",
        "Food Waste in Mess due to lack of real-time meal skip tracking.",
        "Fragmented Student Communication via unverified chat groups.",
        "No accountability or photo proof of work completion by service staff."
    ]
    for b in bullets1:
        p = tf1.add_paragraph()
        p.text = "• " + b
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT_DARK
        p.space_after = Pt(12)

    # Card 2: HostelHub Vision
    card2 = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(1.5), Inches(5.6), Inches(5.3))
    card2.fill.solid()
    card2.fill.fore_color.rgb = CARD_BG
    card2.line.color.rgb = BORDER_GRAY
    tf2 = card2.text_frame
    tf2.word_wrap = True
    p = tf2.paragraphs[0]
    p.text = "🎯 The HostelHub Solution"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = GREEN

    bullets2 = [
        "100% Digital Workflow connecting Students, Wardens, Staff & Admins.",
        "Automated Ticket Dispatching with category-based routing to staff.",
        "Mandatory Image Proof Upload by staff before marking jobs resolved.",
        "Real-Time Roll Call Attendance with instant hostel-wise statistics.",
        "Evening Snacks & Meal Skip System to optimize mess preparation."
    ]
    for b in bullets2:
        p = tf2.add_paragraph()
        p.text = "✔ " + b
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT_DARK
        p.space_after = Pt(12)

    # ==========================================
    # SLIDE 3: SOLUTION OVERVIEW & ARCHITECTURE
    # ==========================================
    slide3 = prs.slides.add_slide(blank_slide_layout)
    add_header(slide3, "Core Platform Architecture & Modules")

    # Left Column Text
    card = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(6.5), Inches(5.3))
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG
    card.line.color.rgb = BORDER_GRAY
    tf = card.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "🚀 Key System Modules"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = DEEP_MAROON

    modules = [
        ("🔧 Ticket Maintenance Engine", "Instant photo complaints, category routing & rating feedback."),
        ("📅 Daily Roll Call Roster", "Binary Present/Absent marking with hostel block filtering."),
        ("☕ Mess Regulations & Snacks", "Weekly meal schedules, evening snacks & skip meal tracking."),
        ("💬 Real-Time Group Chat", "Block & Batch chats powered by Socket.io WebSockets."),
        ("🔒 Profile Security Guard", "Warden re-approval required for Gender or Batch modifications.")
    ]

    for title, desc in modules:
        p = tf.add_paragraph()
        p.text = title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = DEEP_MAROON
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(11.5)
        p2.font.color.rgb = TEXT_MUTED
        p2.space_after = Pt(8)

    # Right Column: Real App Photograph
    if os.path.exists(img_profile):
        slide3.shapes.add_picture(img_profile, Inches(7.6), Inches(1.5), height=Inches(5.3))

    # ==========================================
    # SLIDE 4: STUDENT PORTAL & COMPLAINT RAISING
    # ==========================================
    slide4 = prs.slides.add_slide(blank_slide_layout)
    add_header(slide4, "Student Portal - Instant Issue Raising")

    # Left Box
    card = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(6.5), Inches(5.3))
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG
    card.line.color.rgb = BORDER_GRAY
    tf = card.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "📝 Frictionless Complaint Dispatch"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = DEEP_MAROON

    features = [
        "Categorized Issue Selection: Electrical, Plumbing, Carpentry, Cleaning, Other.",
        "Priority Assignment: Low, Medium, High, Urgent level tags.",
        "Native Photo Capture: Direct camera attachment for physical damage proof.",
        "Real-Time Progress Tracking: 4-stage visual timeline (Ticket Raised -> Staff Assignment -> Work Progress -> Resolution).",
        "Warden Contact Cards: One-tap call/email buttons to reach assigned block wardens instantly."
    ]

    for f in features:
        p = tf.add_paragraph()
        p.text = "• " + f
        p.font.size = Pt(12.5)
        p.font.color.rgb = TEXT_DARK
        p.space_after = Pt(10)

    if os.path.exists(img_complaint):
        slide4.shapes.add_picture(img_complaint, Inches(7.6), Inches(1.5), height=Inches(5.3))

    # ==========================================
    # SLIDE 5: MAINTENANCE WORKFLOW & PROOF
    # ==========================================
    slide5 = prs.slides.add_slide(blank_slide_layout)
    add_header(slide5, "End-to-End Ticket Lifecycle & Proof Verification")

    # 4 Steps Cards Across Widescreen
    steps = [
        ("1. Raise Ticket", "Student logs issue with title, category, priority & issue photo.", GOLD),
        ("2. Assign Staff", "Warden views workload & assigns specialized electrician/plumber.", DEEP_MAROON),
        ("3. Work Progress", "Staff marks job 'In Progress' & visits student room.", CRIMSON),
        ("4. Proof & Feedback", "Staff uploads resolution photo proof. Student rates job 1-5 stars.", GREEN)
    ]

    for i, (stitle, sdesc, color) in enumerate(steps):
        left_pos = Inches(0.8 + i * 2.95)
        box = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, Inches(1.5), Inches(2.75), Inches(5.3))
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_BG
        box.line.color.rgb = BORDER_GRAY
        tf = box.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = stitle
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = color
        p.space_after = Pt(10)

        p2 = tf.add_paragraph()
        p2.text = sdesc
        p2.font.size = Pt(12)
        p2.font.color.rgb = TEXT_DARK

    # ==========================================
    # SLIDE 6: DAILY ROLL CALL ATTENDANCE
    # ==========================================
    slide6 = prs.slides.add_slide(blank_slide_layout)
    add_header(slide6, "Warden Portal - Daily Roll Call Attendance")

    # Left Column Details
    card = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(6.5), Inches(5.3))
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG
    card.line.color.rgb = BORDER_GRAY
    tf = card.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "📅 Streamlined Roll Call Roster"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = DEEP_MAROON

    att_points = [
        "Binary Status Model: Strictly tracks Present & Absent (Outing clutter eliminated).",
        "Hostel Block Filtering: Filter by All Hostels, Boys Hostel 1/2, Girls Hostel 1/2.",
        "Live Summary Metrics Bar: Displays Total Students, Present Count, and Absent Count in real-time.",
        "Bulk Action Controls: One-tap 'Mark All Present' with optional student remark entries.",
        "Student Attendance Metrics: Students view overall attendance % & history."
    ]

    for ap in att_points:
        p = tf.add_paragraph()
        p.text = "✔ " + ap
        p.font.size = Pt(12.5)
        p.font.color.rgb = TEXT_DARK
        p.space_after = Pt(10)

    if os.path.exists(img_attendance):
        slide6.shapes.add_picture(img_attendance, Inches(7.6), Inches(1.5), height=Inches(5.3))

    # ==========================================
    # SLIDE 7: MESS REGULATIONS & MEAL SKIPPING
    # ==========================================
    slide7 = prs.slides.add_slide(blank_slide_layout)
    add_header(slide7, "Mess Regulations, Evening Snacks & Meal Skipping")

    card1 = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.3))
    card1.fill.solid()
    card1.fill.fore_color.rgb = CARD_BG
    card1.line.color.rgb = BORDER_GRAY
    tf1 = card1.text_frame
    tf1.word_wrap = True
    p = tf1.paragraphs[0]
    p.text = "☕ 4-Meal Daily Schedule"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = DEEP_MAROON

    m_points = [
        "Breakfast: Morning healthy options & tea/coffee.",
        "Lunch: Full nutritional meal schedule.",
        "Evening Snacks (☕ Added): Tea, snacks & evening refreshment item.",
        "Dinner: Evening meal schedule with special weekend menus.",
        "Warden Menu Control: Wardens update weekly meal items in real-time."
    ]
    for mp in m_points:
        p = tf1.add_paragraph()
        p.text = "• " + mp
        p.font.size = Pt(12.5)
        p.font.color.rgb = TEXT_DARK
        p.space_after = Pt(8)

    card2 = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(1.5), Inches(5.6), Inches(5.3))
    card2.fill.solid()
    card2.fill.fore_color.rgb = CARD_BG
    card2.line.color.rgb = BORDER_GRAY
    tf2 = card2.text_frame
    tf2.word_wrap = True
    p = tf2.paragraphs[0]
    p.text = "♻️ Meal Skip & Waste Reduction"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = GREEN

    s_points = [
        "One-Tap Skip Toggle: Students mark meals they will skip ahead of time.",
        "Mess Staff Analytics: Live count of meals to prepare, cutting food waste by up to 30%.",
        "Mess Feedback Rating: Students rate food quality, hygiene & quantity.",
        "Warden Feedback Overview: Wardens monitor mess ratings & student comments."
    ]
    for sp in s_points:
        p = tf2.add_paragraph()
        p.text = "✔ " + sp
        p.font.size = Pt(12.5)
        p.font.color.rgb = TEXT_DARK
        p.space_after = Pt(10)

    # ==========================================
    # SLIDE 8: REAL-TIME COMMUNICATION & CHAT
    # ==========================================
    slide8 = prs.slides.add_slide(blank_slide_layout)
    add_header(slide8, "Real-Time Communication & Group Chat")

    card = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(6.5), Inches(5.3))
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG
    card.line.color.rgb = BORDER_GRAY
    tf = card.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "💬 Powered by Socket.io WebSockets"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = DEEP_MAROON

    chat_points = [
        "Automatic Group Membership: Joined into Block Group Chat & Academic Batch Chat on signup.",
        "Warden Group Hub: Wardens broadcast notices to all or specific hostel blocks.",
        "Photo Sharing & Attachment Preview: Share room issues & study notices.",
        "Single-Screen Immersive UI: Bottom navigation stays pinned during scrolling.",
        "Developer Team Cards: Integrated GitHub, LinkedIn, Twitter, Email & 📸 Instagram buttons."
    ]
    for cp in chat_points:
        p = tf.add_paragraph()
        p.text = "• " + cp
        p.font.size = Pt(12.5)
        p.font.color.rgb = TEXT_DARK
        p.space_after = Pt(10)

    if os.path.exists(img_devs):
        slide8.shapes.add_picture(img_devs, Inches(7.6), Inches(1.5), height=Inches(5.3))

    # ==========================================
    # SLIDE 9: TECH STACK & SECURITY GUARDS
    # ==========================================
    slide9 = prs.slides.add_slide(blank_slide_layout)
    add_header(slide9, "Architecture & Critical Security Guards")

    card1 = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.3))
    card1.fill.solid()
    card1.fill.fore_color.rgb = CARD_BG
    card1.line.color.rgb = BORDER_GRAY
    tf1 = card1.text_frame
    tf1.word_wrap = True
    p = tf1.paragraphs[0]
    p.text = "⚙️ Full Stack Architecture"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = DEEP_MAROON

    tech_list = [
        "Frontend: Angular Standalone Components, Vanilla CSS Tokens, RxJS.",
        "Mobile APK: Capacitor Native Android Bridge & Back-Button Guard.",
        "Backend: Node.js, Express.js REST API layer.",
        "Database: PostgreSQL Database with Sequelize ORM & indexing.",
        "Cloud Storage: Supabase Image Storage with Base64 fallback."
    ]
    for t in tech_list:
        p = tf1.add_paragraph()
        p.text = "▪ " + t
        p.font.size = Pt(12.5)
        p.font.color.rgb = TEXT_DARK
        p.space_after = Pt(8)

    card2 = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(1.5), Inches(5.6), Inches(5.3))
    card2.fill.solid()
    card2.fill.fore_color.rgb = CARD_BG
    card2.line.color.rgb = BORDER_GRAY
    tf2 = card2.text_frame
    tf2.word_wrap = True
    p = tf2.paragraphs[0]
    p.text = "🔒 Security Enforcements"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = GOLD

    sec_list = [
        "Profile Re-Approval Rule: Standard profile updates save instantly. Editing Gender or Batch triggers critical warning & requires Warden Re-approval.",
        "Dual-Press Back Button Exit Guard: Back button navigates to Home tab first; second press opens sleek Exit Confirmation Modal.",
        "Role-Based JWT Authorization: Strictly enforces route access for Students, Wardens, Staff & Admins."
    ]
    for s in sec_list:
        p = tf2.add_paragraph()
        p.text = "🛡️ " + s
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_DARK
        p.space_after = Pt(10)

    # ==========================================
    # SLIDE 10: CONCLUSION & THANK YOU
    # ==========================================
    slide10 = prs.slides.add_slide(blank_slide_layout)
    bg10 = slide10.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg10.fill.solid()
    bg10.fill.fore_color.rgb = DEEP_MAROON
    bg10.line.fill.background()

    # Content Box
    tx10 = slide10.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(11.333), Inches(5.0))
    tf10 = tx10.text_frame
    p = tf10.paragraphs[0]
    p.text = "Transforming Hostel Operations with HostelHub"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE

    p2 = tf10.add_paragraph()
    p2.text = "✓ Zero Paperwork  |  ✓ 100% Ticket Accountability  |  ✓ Reduced Mess Waste  |  ✓ Instant Emergency Alerting"
    p2.font.size = Pt(16)
    p2.font.bold = True
    p2.font.color.rgb = GOLD
    p2.space_after = Pt(30)

    p3 = tf10.add_paragraph()
    p3.text = "Thank You!"
    p3.font.size = Pt(44)
    p3.font.bold = True
    p3.font.color.rgb = WHITE

    p4 = tf10.add_paragraph()
    p4.text = "For inquiries, app deployment, or feedback, connect with the HostelHub Engineering Team."
    p4.font.size = Pt(14)
    p4.font.color.rgb = RGBColor(226, 232, 240)

    output_path = r"C:\Users\abhin\Desktop\HostelHub\HostelHub_Presentation.pptx"
    prs.save(output_path)
    print(f"PowerPoint Presentation generated successfully at: {output_path}")

if __name__ == "__main__":
    create_presentation()
